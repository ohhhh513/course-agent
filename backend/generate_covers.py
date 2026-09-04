"""
generate_covers.py - 为 resources 表里的本地资源生成真实封面缩略图。

策略：
  - video : 用 imageio-ffmpeg 自带的 ffmpeg 抽第 1 秒首帧
  - pdf   : 用 poppler 的 pdftoppm 渲染第 1 页
  - ppt   : 用 python-pptx 取第 1 张幻灯片的真实配图；无图则生成占位封面
  - 任何失败都回退到 PIL 生成的带类型/标题占位封面

输出：<project>/assets/resources/covers/{res_id}.jpg  (统一 480x270)
前端按 res_id 直接引用 /assets/resources/covers/{res_id}.jpg，无需改库/重启。
"""
import os
import sys
import shutil
import sqlite3
import subprocess
import tempfile

PROJECT = os.path.dirname(os.path.dirname(os.path.abspath(__file__)))
ASSETS = os.path.join(PROJECT, "assets")
DB = os.path.join(PROJECT, "backend", "app", "data", "course_agent.db")
COVER_DIR = os.path.join(ASSETS, "resources", "covers")
os.makedirs(COVER_DIR, exist_ok=True)

W, H = 480, 270

# ---- 依赖 ----
from PIL import Image, ImageDraw, ImageFont

try:
    import imageio_ffmpeg
    FFMPEG = imageio_ffmpeg.get_ffmpeg_exe()
except Exception:
    FFMPEG = None

PDFTOPPM = shutil.which("pdftoppm")

try:
    from pptx import Presentation
    from pptx.util import Emu
    HAS_PPTX = True
except Exception:
    HAS_PPTX = False


def font(size):
    for cand in ["C:/Windows/Fonts/msyh.ttc", "C:/Windows/Fonts/simhei.ttf",
                 "C:/Windows/Fonts/arial.ttf"]:
        if os.path.exists(cand):
            try:
                return ImageFont.truetype(cand, size)
            except Exception:
                pass
    return ImageFont.load_default()


def local_path(url):
    if not url:
        return None
    rel = url.split("/assets/", 1)[-1]
    return os.path.join(ASSETS, rel)


def cover_fit(img):
    """等比缩放到 WxH 画布居中（cover-fit 简化版：contain + 深底）"""
    bg = Image.new("RGB", (W, H), (24, 28, 38))
    img.thumbnail((W, H), Image.LANCZOS)
    bg.paste(img, ((W - img.width) // 2, (H - img.height) // 2))
    return bg


def placeholder(res_type, title):
    colors = {
        "video": ((196, 30, 58), "▶ 教学视频"),
        "pdf": ((28, 96, 168), "📄 教材文献"),
        "ppt": ((198, 112, 30), "📊 课堂PPT"),
        "doc": ((28, 96, 168), "📄 文档"),
    }
    bgc, label = colors.get(res_type, ((60, 64, 72), "资源"))
    img = Image.new("RGB", (W, H), bgc)
    d = ImageDraw.Draw(img)
    # 顶部色条
    d.rectangle([0, 0, W, 6], fill=(255, 255, 255))
    f_l = font(26)
    f_t = font(22)
    d.text((24, 96), label, fill=(255, 255, 255), font=f_l)
    t = (title or "")[:18]
    d.text((24, 150), t, fill=(235, 238, 245), font=f_t)
    return img


def save_cover(res_id, img):
    out = os.path.join(COVER_DIR, f"{res_id}.jpg")
    img.convert("RGB").save(out, "JPEG", quality=85)
    return out


def extract_video(src, res_id, title):
    if not FFMPEG or not os.path.exists(src):
        return None
    tmp = os.path.join(tempfile.gettempdir(), f"{res_id}_frame.png")
    try:
        subprocess.run(
            [FFMPEG, "-y", "-ss", "1", "-i", src, "-frames:v", "1",
             "-vf", "scale=480:-1", "-q:v", "2", tmp],
            stdout=subprocess.DEVNULL, stderr=subprocess.DEVNULL, timeout=60,
        )
        if os.path.exists(tmp) and os.path.getsize(tmp) > 0:
            return cover_fit(Image.open(tmp))
    except Exception:
        pass
    finally:
        try:
            os.remove(tmp)
        except Exception:
            pass
    return None


def extract_pdf(src, res_id, title):
    if not PDFTOPPM or not os.path.exists(src):
        return None
    try:
        with tempfile.TemporaryDirectory() as td:
            base = os.path.join(td, "p")
            subprocess.run(
                [PDFTOPPM, "-png", "-f", "1", "-l", "1", "-r", "50", src, base],
                stdout=subprocess.DEVNULL, stderr=subprocess.DEVNULL, timeout=60,
            )
            png = base + "-1.png"
            if os.path.exists(png):
                return cover_fit(Image.open(png))
    except Exception:
        pass
    return None


def extract_ppt(src, res_id, title):
    if not HAS_PPTX or not os.path.exists(src):
        return None
    try:
        prs = Presentation(src)
        slide = prs.slides[0]
        for shape in slide.shapes:
            if shape.shape_type == 13:  # PICTURE
                img = shape.image
                from io import BytesIO
                return cover_fit(Image.open(BytesIO(img.blob)))
    except Exception:
        pass
    return None


def main():
    con = sqlite3.connect(DB)
    cur = con.cursor()
    cur.execute("SELECT res_id, type, title, url FROM resources")
    rows = cur.fetchall()
    con.close()

    ok = 0
    for res_id, rtype, title, url in rows:
        src = local_path(url)
        img = None
        if rtype == "video":
            img = extract_video(src, res_id, title)
        elif rtype in ("pdf", "doc"):
            img = extract_pdf(src, res_id, title)
        elif rtype == "ppt":
            img = extract_ppt(src, res_id, title)
        if img is None:
            img = placeholder(rtype, title)
        save_cover(res_id, img)
        ok += 1
        print(f"  cover {res_id} ({rtype}) <- {os.path.basename(src) if src else 'placeholder'}")

    print(f"DONE: 生成 {ok} 张封面 -> {COVER_DIR}")


if __name__ == "__main__":
    main()
